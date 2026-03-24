#!/usr/bin/env python3
"""Generate CPS-IDS presentation as .pptx (importable to Google Slides)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colours
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0xE9, 0x45, 0x60)
BLUE = RGBColor(0x1E, 0x40, 0xAF)
AMBER = RGBColor(0x92, 0x40, 0x0E)
GREEN = RGBColor(0x06, 0x5F, 0x46)
GREY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── helpers ──────────────────────────────────────────────────────────────────

def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT,
             space_before=Pt(4), space_after=Pt(2), level=0, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.level = level
    p.space_before = space_before
    p.space_after = space_after
    return p

def add_title(slide, text, subtitle=None):
    add_textbox(slide, 0.5, 0.3, 12.3, 0.8, text, font_size=36, bold=True, color=DARK)
    if subtitle:
        add_textbox(slide, 0.5, 1.05, 12.3, 0.5, subtitle, font_size=18, color=GREY)

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_stat_card(slide, left, top, number, label, accent=DARK):
    add_rect(slide, left, top, 2.4, 1.5, LIGHT_BG)
    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                  Inches(0.08), Inches(1.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    # number
    add_textbox(slide, left + 0.2, top + 0.15, 2.1, 0.7, str(number),
                font_size=32, bold=True, color=accent, alignment=PP_ALIGN.CENTER)
    # label
    add_textbox(slide, left + 0.2, top + 0.85, 2.1, 0.6, label,
                font_size=11, color=GREY, alignment=PP_ALIGN.CENTER)

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        p.level = 0
    return tf

def add_table(slide, left, top, width, height, rows, col_widths=None):
    """Add a table. rows[0] is the header row."""
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r, row_data in enumerate(rows):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    return table


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════════════════════
s = add_blank_slide()
# dark banner
banner = slide_shapes = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(13.333), Inches(7.5))
banner.fill.solid()
banner.fill.fore_color.rgb = DARK
banner.line.fill.background()

add_textbox(s, 1, 2.0, 11.3, 1.2, "CPS-IDS", font_size=64, bold=True, color=WHITE,
            alignment=PP_ALIGN.CENTER)
add_textbox(s, 1, 3.3, 11.3, 0.8,
            "AI/ML-Driven Intrusion Detection for Cyber-Physical Systems",
            font_size=28, color=RGBColor(0xBB, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)
add_textbox(s, 1, 5.0, 11.3, 0.5, "CITY3116 \u2014 Advanced Computer Forensics  |  2025\u201326",
            font_size=16, color=GREY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: CPS Architecture
# ═══════════════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_title(s, "CPS Testbed Architecture", "Physical Arduino plant + software simulation (Purdue model)")

# Architecture diagram as text box
arch_text = (
    "Level 3 \u2014 SCADA PC / IDS / Attacker\n"
    "  \u2502  Ethernet (Modbus TCP)\n"
    "Level 2 \u2014 TP-Link LS1005G Switch\n"
    "  \u2502\n"
    "  \u251c\u2500 PLC 1 (Mega + W5500)  192.168.1.20\n"
    "  \u2502    \u2502  RS-485 (Modbus RTU)\n"
    "  \u2514\u2500 PLC 2 (Uno + MAX485)\n"
    "  \u2502\n"
    "Level 1 \u2014 Control Logic\n"
    "  \u2502\n"
    "Level 0 \u2014 Sensors & Actuators"
)
rect = add_rect(s, 0.5, 1.7, 5.5, 4.5, LIGHT_BG)
add_textbox(s, 0.8, 1.9, 5.0, 4.2, arch_text, font_size=15, color=DARK, font_name="Consolas")

add_bullet_list(s, 6.5, 1.7, 6.3, 5.0, [
    "\u2022  PLC 1 (Arduino Mega): Modbus TCP server, pump/valve control, RFID auth, LCD HMI",
    "\u2022  PLC 2 (Arduino Uno): Modbus RTU slave, 5 sensors (ultrasonic, temp, sound, motion, water)",
    "\u2022  Software twin: MiniCPS/pymodbus simulation with identical register map and control logic",
    "\u2022  Plant dashboard: FastAPI + WebSocket GUI with live SVG visualisation and attack terminal",
    "\u2022  Control logic: pump ON < 30%, OFF > 85%; valve synced; temp alarm > 40\u00b0C",
], font_size=16)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: IDS Design
# ═══════════════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_title(s, "Multi-Layer IDS Design", "Implemented in Rust (6-crate workspace) for memory safety and performance")

add_stat_card(s, 1.0, 2.0, "20", "Rule Engine Rules\nSYN flood, port scan,\nModbus abuse, injection", BLUE)
add_stat_card(s, 4.0, 2.0, "2", "Modbus Analysers\nWrite-rate anomaly\nRead flood detection", AMBER)
add_stat_card(s, 7.0, 2.0, "3", "ML Models\nRandom Forest, CNN+LSTM\nIsolation Forest", GREEN)
add_stat_card(s, 10.0, 2.0, "78", "Flow Features\nExtracted per flow for\nCNN+LSTM inference", ACCENT)

add_textbox(s, 0.5, 4.2, 12.3, 0.5,
            "Detection Pipeline:   Rule Engine  \u2192  Modbus Analysis  \u2192  CNN+LSTM Inference",
            font_size=20, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)

add_bullet_list(s, 0.8, 5.0, 5.8, 2.3, [
    "\u2022  ids-common: shared types, config, errors",
    "\u2022  ids-collector: packet capture, flow tracking",
    "\u2022  ids-preprocess: dataset loaders, SMOTE, scaling",
], font_size=14, color=GREY)

add_bullet_list(s, 6.8, 5.0, 5.8, 2.3, [
    "\u2022  ids-engine: models, rules, ML inference, binaries",
    "\u2022  ids-response: alerter, blocker, SIEM export",
    "\u2022  ids-dashboard: web UI stub (axum)",
], font_size=14, color=GREY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: ML Models
# ═══════════════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_title(s, "AI/ML Model Architecture", "3 complementary models trained on 3 public datasets")

# Left: CNN+LSTM
add_textbox(s, 0.5, 1.6, 6.0, 0.4, "CNN+LSTM (~297K parameters)", font_size=20, bold=True, color=DARK)
add_bullet_list(s, 0.5, 2.1, 6.0, 2.5, [
    "\u2022  Conv1d \u2192 ReLU \u2192 BatchNorm \u2192 Conv1d (feature extraction)",
    "\u2022  LSTM: 2 layers, 128 hidden units (temporal patterns)",
    "\u2022  FC: 128 \u2192 64 \u2192 5 classes, dropout 0.3",
    "\u2022  PyTorch training \u2192 ONNX export for live inference",
    "\u2022  Early stopping with patience=5, LR decay on plateau",
], font_size=15)

add_textbox(s, 0.5, 4.3, 6.0, 0.4, "Ensemble (RF + Isolation Forest)", font_size=20, bold=True, color=DARK)
add_bullet_list(s, 0.5, 4.8, 6.0, 1.8, [
    "\u2022  Random Forest: 100 trees, max depth 20 (supervised)",
    "\u2022  Isolation Forest: anomaly detection for zero-days (unsupervised)",
    "\u2022  Hard voting: 70% RF / 30% IForest weighted blend",
], font_size=15)

# Right: 5-class scheme + datasets
add_table(s, 7.0, 1.6, 5.8, 2.0, [
    ["Class", "Category", "Example Attacks"],
    ["0", "Normal", "Benign traffic"],
    ["1", "DoS", "DDoS, SYN flood, Slowloris"],
    ["2", "Probe", "Port scan, reconnaissance"],
    ["3", "R2L", "Brute force, web attacks"],
    ["4", "U2R", "Buffer overflow, rootkit"],
], col_widths=[0.6, 1.2, 4.0])

add_textbox(s, 7.0, 4.3, 5.8, 0.4, "Training Datasets", font_size=20, bold=True, color=DARK)
add_table(s, 7.0, 4.8, 5.8, 1.8, [
    ["Model", "Dataset", "Year", "Features", "SMOTE"],
    ["A", "NSL-KDD", "1999", "122", "Yes"],
    ["B", "CIC-IDS2017", "2017", "78", "No"],
    ["C", "UNSW-NB15", "2015", "76", "Yes"],
    ["D", "Combined (A+B+C)", "Mixed", "276", "No"],
], col_widths=[0.8, 2.0, 0.8, 1.1, 1.1])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Training Results
# ═══════════════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_title(s, "Training Results")

add_table(s, 0.5, 1.5, 12.3, 2.5, [
    ["Model", "Dataset", "Features", "Random Forest", "CNN+LSTM", "Isolation Forest", "RF+IForest", "FPR"],
    ["A", "NSL-KDD", "122", "67.26%*", "65.94%*", "77.24%", "67.26%*", "0.0301"],
    ["B", "CIC-IDS2017", "78", "99.73%", "99.83%", "81.37%", "99.73%", "0.0006"],
    ["C", "UNSW-NB15", "76", "93.90%", "92.95%", "82.50%", "93.84%", "0.0251"],
    ["D", "Combined", "276", "97.80%", "97.67%", "80.20%", "97.82%", "0.0043"],
], col_widths=[0.8, 1.8, 1.0, 1.6, 1.5, 1.7, 1.5, 1.0])

add_textbox(s, 0.5, 4.0, 12.3, 0.4,
            "*Model A limited by NSL-KDD binary-only test labels (Probe/R2L/U2R collapsed)",
            font_size=12, color=GREY)

add_stat_card(s, 1.5, 4.7, "99.83%", "Best Accuracy\nCNN+LSTM, Model B", ACCENT)
add_stat_card(s, 4.8, 4.7, "0.06%", "False Positive Rate\nModel B", ACCENT)
add_stat_card(s, 8.1, 4.7, "97.82%", "Cross-Era Generalisation\nEnsemble, Model D", ACCENT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Live Monitor
# ═══════════════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_title(s, "Live IDS Monitor", "Real-time packet capture, flow analysis, and ML inference")

add_textbox(s, 0.5, 1.6, 6.0, 0.4, "Architecture", font_size=20, bold=True)
add_bullet_list(s, 0.5, 2.1, 6.0, 2.5, [
    "\u2022  Rust binary captures packets via pnet (raw sockets)",
    "\u2022  Bidirectional flow tracking with configurable timeout",
    "\u2022  Extracts 78 CIC-IDS2017 features per completed flow",
    "\u2022  Python subprocess runs ONNX model via JSON Lines IPC",
    "\u2022  Time-windowed inference every 10s on active flows",
], font_size=15)

add_textbox(s, 0.5, 4.3, 6.0, 0.4, "Modbus-Specific Detection", font_size=20, bold=True)
add_bullet_list(s, 0.5, 4.8, 6.0, 1.8, [
    "\u2022  Write-rate anomaly: >2.0 writes/sec triggers alert",
    "\u2022  Read flood detection: >50 reads/sec flagged as DoS",
    "\u2022  Domain-adapted classifier for Modbus flow patterns",
], font_size=15)

add_stat_card(s, 7.5, 2.0, "~1ms", "ML Inference\nLatency per prediction", GREEN)
add_stat_card(s, 10.3, 2.0, "10s", "Inference Window\nActive flow analysis", BLUE)

add_textbox(s, 7.0, 4.3, 5.8, 0.4, "Alert Output", font_size=20, bold=True)
add_bullet_list(s, 7.0, 4.8, 5.8, 2.0, [
    "\u2022  Structured JSON (alerts.jsonl) with severity levels",
    "\u2022  Per-alert: category, confidence, model source, 5-tuple",
    "\u2022  SIEM-ready: CEF/syslog export capability",
    "\u2022  Graceful shutdown with final flow expiry and stats",
], font_size=15)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: Attack Scenario
# ═══════════════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_title(s, "Forensic Investigation: Attack Scenario",
          "7 attacks executed via automated pipeline (tmux + tcpdump + IDS + attack runner)")

add_table(s, 0.5, 1.7, 12.3, 4.0, [
    ["#", "Attack", "Modbus FC", "Duration", "Description"],
    ["1", "Command Injection", "FC 0x05 (Write Coil)", "45s", "Forces pump ON/OFF every 5 seconds"],
    ["2", "Pump Oscillation", "FC 0x05 (Write Coil)", "45s", "Rapid pump toggle every 2s (Stuxnet-style)"],
    ["3", "Valve Manipulation", "FC 0x06 (Write Reg)", "45s", "Random valve positions (0\u2013180\u00b0) every 3s"],
    ["4", "Replay Attack", "FC 0x05 (Write Coil)", "45s", "Replays recorded ON/OFF command sequence"],
    ["5", "Sensor Spoofing", "FC 0x06 (Write Reg)", "45s", "Writes fake tank levels (0\u2013100) every 2s"],
    ["6", "Modbus Flood (DoS)", "FC 0x03 (Read Reg)", "20s", "100\u00d7 read requests per iteration"],
    ["7", "Multi-Stage", "Mixed", "30s", "Recon (10s) \u2192 manipulation (20s) \u2192 flood"],
], col_widths=[0.5, 2.2, 2.5, 1.1, 6.0])

add_textbox(s, 0.5, 5.9, 12.3, 0.8,
            "Evidence captured:  full-session.pcap (tcpdump)  |  alerts.jsonl (IDS)  |  attack-manifest.json (timestamps)",
            font_size=15, color=GREY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Investigation Results
# ═══════════════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_title(s, "Investigation Results")

add_stat_card(s, 0.8, 1.7, "17", "ML Alerts\nAcross all 7 attacks", GREEN)
add_stat_card(s, 3.5, 1.7, "300", "Write-Rate Alerts\nModbus anomaly detection", AMBER)
add_stat_card(s, 6.2, 1.7, "8,126", "Flood Detection\nDoS & multi-stage attacks", BLUE)

add_textbox(s, 9.5, 1.7, 3.3, 0.4, "Detection Coverage", font_size=18, bold=True)
add_table(s, 9.5, 2.2, 3.3, 3.5, [
    ["Attack", "Rule", "Modbus", "ML"],
    ["Cmd Injection", "\u2713", "\u2713", "\u2713"],
    ["Pump Oscillation", "\u2713", "\u2713", "\u2713"],
    ["Valve Manip.", "\u2713", "\u2713", "\u2713"],
    ["Replay", "\u2713", "\u2713", "\u2713"],
    ["Sensor Spoof", "\u2713", "\u2713", "\u2713"],
    ["Modbus Flood", "\u2713", "\u2713", "\u2713"],
    ["Multi-Stage", "\u2713", "\u2713", "\u2713"],
], col_widths=[1.3, 0.6, 0.8, 0.6])

add_textbox(s, 0.5, 3.7, 8.5, 0.4, "Forensic Tools & Evidence", font_size=20, bold=True)
add_bullet_list(s, 0.5, 4.2, 8.5, 2.8, [
    "\u2022  Wireshark: pcap analysis with Modbus TCP dissector (modbus.func_code filter)",
    "\u2022  tcpdump: live packet capture on loopback port 5502 during attack sequence",
    "\u2022  Custom IDS monitor: 3-layer real-time alert generation with JSON structured output",
    "\u2022  Attack manifest: JSON timestamps for precise correlation of alerts to attack phases",
    "\u2022  Key finding: multi-layer detection achieves 100% attack coverage \u2014 every attack type detected by all 3 layers",
], font_size=15)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: Evaluation & Recommendations
# ═══════════════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_title(s, "Evaluation & Recommendations")

add_textbox(s, 0.5, 1.6, 6.0, 0.4, "Strengths", font_size=20, bold=True, color=GREEN)
add_bullet_list(s, 0.5, 2.1, 6.0, 2.0, [
    "\u2022  CNN+LSTM achieves 99.83% accuracy with 0.06% FPR",
    "\u2022  Multi-layer defence-in-depth (rules + rate + ML)",
    "\u2022  Cross-dataset generalisation at 97.8% (Model D)",
    "\u2022  Real-time inference <1ms via ONNX Runtime",
    "\u2022  Ensemble improves minority class recall (R2L +7%)",
], font_size=15)

add_textbox(s, 0.5, 4.0, 6.0, 0.4, "Limitations", font_size=20, bold=True, color=ACCENT)
add_bullet_list(s, 0.5, 4.5, 6.0, 2.0, [
    "\u2022  U2R class hardest (extremely rare training samples)",
    "\u2022  Model A degraded by NSL-KDD binary test labels",
    "\u2022  Domain gap: training data \u2260 live Modbus traffic",
    "\u2022  Rule engine generates false positives on localhost",
], font_size=15)

add_textbox(s, 7.0, 1.6, 5.8, 0.4, "Recommendations", font_size=20, bold=True, color=BLUE)
add_bullet_list(s, 7.0, 2.1, 5.8, 2.0, [
    "\u2022  Domain adaptation: fine-tune on captured Modbus traffic",
    "\u2022  Federated learning: train across CPS sites privately",
    "\u2022  Adaptive thresholds: tune rule engine per deployment",
    "\u2022  SIEM integration: CEF/syslog export for SOC workflows",
], font_size=15)

add_textbox(s, 7.0, 4.0, 5.8, 0.4, "Standards Alignment", font_size=20, bold=True, color=DARK)
add_bullet_list(s, 7.0, 4.5, 5.8, 2.0, [
    "\u2022  ISO 27001: A.12.4 logging, A.16 incident management",
    "\u2022  IEC 62443: defence-in-depth for industrial control systems",
    "\u2022  NIST CSF: Detect & Respond functions",
    "\u2022  OWASP Top 10: injection prevention, access control",
], font_size=15)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: Summary & Q&A
# ═══════════════════════════════════════════════════════════════════════════════
s = add_blank_slide()
add_title(s, "Summary")

add_stat_card(s, 1.0, 1.7, "6", "Rust Crates", DARK)
add_stat_card(s, 3.7, 1.7, "4", "Model Variants", DARK)
add_stat_card(s, 6.4, 1.7, "99.83%", "Best Accuracy", ACCENT)
add_stat_card(s, 9.1, 1.7, "7 / 7", "Attacks Detected", GREEN)

add_bullet_list(s, 0.8, 3.8, 11.7, 2.8, [
    "\u2022  Built a complete CPS testbed (Arduino hardware + MiniCPS software simulation)",
    "\u2022  Trained CNN+LSTM, Random Forest, and Isolation Forest on 3 public IDS datasets",
    "\u2022  Deployed live IDS monitor with real-time ONNX inference (~1ms per prediction)",
    "\u2022  Conducted forensic investigation detecting all 7 attack types across 3 detection layers",
    "\u2022  Multi-layer approach compensates for individual model weaknesses \u2014 100% attack coverage",
], font_size=18)

# Thank you
add_textbox(s, 0.5, 6.0, 12.3, 0.8, "Thank you \u2014 Questions?",
            font_size=32, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════════
output_path = "/home/vt/Documents/BSC/advanced-computer-forensics/CITY3116-1/presentation/CPS-IDS-Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
